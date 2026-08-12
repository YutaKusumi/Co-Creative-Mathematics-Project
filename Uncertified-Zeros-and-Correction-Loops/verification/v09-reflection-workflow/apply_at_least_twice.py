# -*- coding: utf-8 -*-
"""「二回きく」→「二回以上きく」の全面適用（登録者裁定・2026-08-11・案(あ)）。

理由（改訂記録に残す）: 第4条の根拠は ε<3/n であり、n が増えるほど上界は下がる。
「二回きく」という標語は「二回で足りる」という偽の十分性を含みうるのに対し、
「二回以上きく」は本稿の論理と一致する。やさしい日本語版の起草中に登録者が検出した。

対象: 原典 JA/EN（§7.1 第4条・付録D(一)④）＋教材5文書＋pptx＋各 README。
実行: proposals で python apply_at_least_twice.py
"""
import io, hashlib, copy
from pptx import Presentation

R = 'C:/Users/PC/Desktop/GitHub-Repositories/Co-Creative-Mathematics-Project/Uncertified-Zeros-and-Correction-Loops/'
sha = lambda p: hashlib.sha256(io.open(p, 'rb').read().replace(b'\r\n', b'\n')).hexdigest()[:16].upper()
def fix(p, pairs):
    s = io.open(R + p, encoding='utf-8').read()
    done = 0
    for old, new, l in pairs:
        n = s.count(old)
        if new in s:      # 適用済み（再実行できるようにする。新が旧を含む場合も飛ばす——初版は二重書き込みを起こした）
            done += 1; continue
        assert n == 1, 'アンカー %s / %s: %d件' % (p.split('/')[-1], l, n)
        s = s.replace(old, new)
    io.open(R + p, 'w', encoding='utf-8', newline='').write(s)
    print('  %-46s %s%s' % (p.split('/')[-1], sha(R + p), '  （うち適用済み%d件）' % done if done else ''))

print('== 原典 JA ==')
fix('JA/uncertified-zeros-and-correction-loops-JA.md', [
 ('**版**: 第一草稿・完成版 v0.9.7（2026年8月11日）', '**版**: 第一草稿・完成版 v0.9.8（2026年8月11日）', '版'),
 ('だから、大事なことは二回きく。', 'だから、大事なことは二回以上きく。', '§7.1第4条'),
 ('④大事なことは、二回きく。別の日に、できれば別の相手にも。', '④大事なことは、二回以上きく。別の日に、できれば別の相手にも。', '付録D(一)④'),
 ('やさしい日本語としての妥当性は当事者にも専門家にも検分されていない）を付録G に明記（本文の主張には触れない改訂）。',
  '''やさしい日本語としての妥当性は当事者にも専門家にも検分されていない）を付録G に明記（本文の主張には触れない改訂）。 v0.9.8 第7.1節 第4条と付録D(一)④の「二回きく」を「**二回以上きく**」に改めた——第4条の根拠は ε<3/n であり、n が増えるほど上界は下がる。「二回きく」という標語は「二回で足りる」という偽の十分性を含みうるのに対し、「二回以上きく」は本稿の論理と一致する（主張を弱める変更ではなく、標語を論理に合わせる変更である）。付録G(3) のやさしい日本語版の起草中に登録者が検出し、派生教材の全箇所へ同時に反映した（登録者裁定・2026年8月11日）。''', '改訂記録'),
])

print('== 原典 EN ==')
fix('EN/uncertified-zeros-and-correction-loops-EN.md', [
 ('**Version**: First Draft, Completed Version v0.9.7 (August 11, 2026)',
  '**Version**: First Draft, Completed Version v0.9.8 (August 11, 2026)', '版'),
 ('So ask what matters twice.', 'So ask what matters at least twice.', '§7.1第4条'),
 ('④ Ask about anything important twice — on a different day, and if possible, of a different party too.',
  '④ Ask about anything important at least twice — on a different day, and if possible, of a different party too.', '付録D(1)④'),
 ('are stated in Appendix G (a revision that does not touch the claims of the body text).',
  '''are stated in Appendix G (a revision that does not touch the claims of the body text). v0.9.8 "Ask what matters twice" in §7.1 Article 4 and in Appendix D(1)④ was changed to "**ask what matters at least twice**" — the ground of Article 4 is ε<3/n, and the bound falls as n grows. The maxim "ask twice" can carry a false sufficiency ("twice is enough"), whereas "at least twice" agrees with this paper's own logic (this strengthens rather than weakens the claim; it aligns the maxim with the reasoning). The registrant detected it while the Easy Japanese version of Appendix G(3) was being drafted, and it was reflected simultaneously in every corresponding place in the derivative teaching materials (registrant's adjudication, August 11, 2026).''', '改訂記録'),
])

print('== 教材（md） ==')
fix('teaching-materials/README.md', [
 ('③二回きく ④逆もきく', '③二回以上きく ④逆もきく', '技一覧')])
fix('teaching-materials/ai-relationship-lesson-handout-JA.md', [
 ('**四つの技（たしかめる・中身をきく・二回きく・逆もきく）**', '**四つの技（たしかめる・中身をきく・二回以上きく・逆もきく）**', 'ねらい2'),
 ('- **技③二回きく**: 大事なことは、もう一回きく。', '- **技③二回以上きく**: 大事なことは、二回以上きく。', '技③'),
 ('3. **だいじなことは、二回きく。**（べつの場所でも）', '3. **だいじなことは、二回以上きく。**（べつの場所でも）', 'やくそく3')])
fix('teaching-materials/ai-relationship-lesson-script-JA.md', [
 ('### 技③　二回きく', '### 技③　二回以上きく', '見出し'),
 ('三つ目の技。**大事なことは、二回きく。**', '三つ目の技。**大事なことは、二回以上きく。**', '本文'),
 ('〔板書〕**技③　大事なことは、二回きく**', '〔板書〕**技③　大事なことは、二回以上きく**', '板書'),
 ('3. **だいじなことは、二回きく。**', '3. **だいじなことは、二回以上きく。**', 'やくそく3')])
fix('teaching-materials/ai-relationship-lesson-slide-script-JA.md', [
 ('技③「二回きく」', '技③「二回以上きく」', 'スライド6リスト'),
 ('## スライド9　技③　二回きく', '## スライド9　技③　二回以上きく', 'スライド9見出し'),
 ('技の三つ目。**大事なことは、もう一回きく。**', '技の三つ目。**大事なことは、二回以上きく。**', 'スライド9本文'),
 ('三、**だいじなことは、二回きく。**', '三、**だいじなことは、二回以上きく。**', 'やくそく3')])

print('== 教材（pptx） ==')
P = Presentation(R + 'teaching-materials/ai-relationship-lesson-slides-JA.pptx')
def set_text(shape, lines):
    tf = shape.text_frame
    paras = [p for p in tf.paragraphs]
    src = next(p for p in paras if p.runs)
    for i, line in enumerate(lines):
        rs = paras[i].runs
        rs[0].text = line
        for r in rs[1:]: r._r.getparent().remove(r._r)
n = 0
for i, sl in enumerate(P.slides, 1):
    for sh in sl.shapes:
        if not sh.has_text_frame: continue
        for pa in sh.text_frame.paragraphs:
            t = ''.join(r.text for r in pa.runs)
            if t in ('二回きく', '技③　二回きく', 'だいじなことは、二回きく。', '大事なことは、もう一回 きく。'):
                new = {'二回きく': '二回以上きく', '技③　二回きく': '技③　二回以上きく',
                       'だいじなことは、二回きく。': 'だいじなことは、二回以上きく。',
                       '大事なことは、もう一回 きく。': '大事なことは、二回以上 きく。'}[t]
                pa.runs[0].text = new
                for r in pa.runs[1:]: r._r.getparent().remove(r._r)
                n += 1
                print('  スライド%-2d %s → %s' % (i, t, new))
P.save(R + 'teaching-materials/ai-relationship-lesson-slides-JA.pptx')
print('  pptx %d箇所 / SHA %s' % (n, sha(R + 'teaching-materials/ai-relationship-lesson-slides-JA.pptx')))
