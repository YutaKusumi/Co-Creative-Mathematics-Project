# -*- coding: utf-8 -*-
"""授業スライド pptx の同期（対応表 v1・M1/M2/M7/M13）。
 1. スライド6（全体地図）: 三つの技→四つの技・カードを5枚組に再配置
 2. 新スライド8「技② 中身をきく」を挿入（スライド8を雛形に複製）
 3. 旧8→9「技③ 二回きく」: 表題の番号と、赤枠（同じ答えでも正しいとはかぎらない）を追加
 4. 旧9→10「技④ 逆もきく」: 表題の番号
 5. 旧16→17「六つの やくそく」: 2番目に新項目を挿入し番号を振り直す
 6. ページ番号: 旧8以降を +1
 ※ スライド15（→16）の「技①たしかめるの一番強い形」は技①のまま（繰り下げ対象外）
実行: proposals で python apply_teaching_pptx.py
"""
import copy, sys, io, hashlib
from pptx import Presentation
from pptx.util import Pt, Emu

SRC = 'C:/Users/PC/Downloads/ai-relationship-lesson-slides-JA.pptx'
DST = 'C:/Users/PC/Desktop/Ryokai-OS-Verification/proposals/_teaching/ai-relationship-lesson-slides-JA.pptx'
TEAL, DARK, RED, GRAY2 = '028090', '2B3A42', 'F96167', '5C6E77'

P = Presentation(SRC)
S = P.slides
def sh_of(slide, name):
    for sh in slide.shapes:
        if sh.name == name: return sh
    raise KeyError('%s not found' % name)

def set_para(p, text):
    """段落のテキストを差し替える（最初の run の書式を保つ）。"""
    rs = p.runs
    assert rs, '空段落には書けない'
    rs[0].text = text
    for r in rs[1:]:
        r._r.getparent().remove(r._r)

def set_text(shape, lines):
    """図形のテキストを行リストで差し替える。段落数が足りなければ複製して増やす。"""
    tf = shape.text_frame
    paras = [p for p in tf.paragraphs]
    src = None
    for p in paras:
        if p.runs: src = p; break
    assert src is not None, '%s に run が無い' % shape.name
    while len(paras) < len(lines):
        new = copy.deepcopy(src._p)
        src._p.getparent().append(new)
        paras = [p for p in tf.paragraphs]
    for i, line in enumerate(lines):
        if not paras[i].runs:
            paras[i]._p.getparent().replace(paras[i]._p, copy.deepcopy(src._p))
            paras = [p for p in tf.paragraphs]
        set_para(paras[i], line)
    for p in paras[len(lines):]:
        p._p.getparent().remove(p._p)

def set_color(shape, rgb):
    from pptx.dml.color import RGBColor
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor.from_string(rgb)

def set_size(shape, pt):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(pt)

def place(shape, left=None, top=None, width=None, height=None):
    if left is not None:   shape.left = Pt(left)
    if top is not None:    shape.top = Pt(top)
    if width is not None:  shape.width = Pt(width)
    if height is not None: shape.height = Pt(height)

# ============================================================
# 1. スライド6 — 全体地図を5枚組へ
# ============================================================
s6 = S[5]
set_text(sh_of(s6, 'Text 0'), ['かしこく つかう ―― 四つの技 ＋ 一つのルール'])

# 既存4組（枠・円・絵・技番号・名前）と、複製して作る5組目
GROUPS = [('Shape 1', 'Shape 2', 'Image 0', 'Text 3', 'Text 4'),
          ('Shape 5', 'Shape 6', 'Image 1', 'Text 7', 'Text 8'),
          ('Shape 9', 'Shape 10', 'Image 2', 'Text 11', 'Text 12'),
          ('Shape 13', 'Shape 14', 'Image 3', 'Text 15', 'Text 16')]
tree = s6.shapes._spTree
clones = []
for nm in GROUPS[1]:                       # 技②の組を雛形に5組目を作る（画像の関係IDは同一スライド内で共有）
    el = copy.deepcopy(sh_of(s6, nm)._element)
    tree.append(el)
    clones.append(el)
newnames = ['Shape 90', 'Shape 91', 'Image 90', 'Text 90', 'Text 91']
for el, nn in zip(clones, newnames):
    nv = el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
    if nv is None:
        nv = el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    nv.set('name', nn); nv.set('id', str(900 + newnames.index(nn)))
GROUPS.append(tuple(newnames))

# 5組の新しい配置（左57 / 幅157 / 間19.25 → 右端919）
W, GAP, X0 = 157.0, 19.25, 57.0
LABELS = [('技①', 'たしかめる'), ('技②', '中身をきく'), ('技③', '二回きく'),
          ('技④', '逆もきく'), ('ルール', 'ボタンは自分で')]
order = [GROUPS[0], GROUPS[4], GROUPS[1], GROUPS[2], GROUPS[3]]   # ①／新②／旧②→③／旧③→④／ルール
for i, (g, (lab, nam)) in enumerate(zip(order, LABELS)):
    x = X0 + i * (W + GAP)
    card, circ, img, tlab, tnam = [sh_of(s6, n) for n in g]
    place(card, left=x, width=W)
    place(circ, left=x + (W - 79) / 2)
    place(img,  left=x + (W - 44) / 2)
    place(tlab, left=x, width=W)
    place(tnam, left=x + 7, width=W - 14)
    set_text(tlab, [lab]); set_text(tnam, [nam])
    set_size(tnam, 17)

# ============================================================
# 2. 新スライド8「技② 中身をきく」——旧スライド8を雛形に複製して挿入
# ============================================================
tpl = S[7]                                  # 旧スライド8（技② 二回きく）
new = S.add_slide(tpl.slide_layout)
for sh in list(new.shapes):                 # レイアウト由来のプレースホルダを除去
    sh._element.getparent().remove(sh._element)
img_rid = None
for rid, rel in tpl.part.rels.items():
    if rel.reltype.endswith('/image'):
        img_rid = new.part.relate_to(rel.target_part, rel.reltype)
        old_rid = rid
for sh in tpl.shapes:
    el = copy.deepcopy(sh._element)
    x = el.xml
    new.shapes._spTree.append(el)
if img_rid:                                 # 画像の関係IDを新スライドのものへ張り替え
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
    for blip in new.shapes._spTree.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
        blip.set(ns, img_rid)
# 位置を8番目へ移動（末尾から）
sldIdLst = P.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-1]); sldIdLst.insert(7, ids[-1])

s8 = S[7]                                   # 新スライド8
set_text(sh_of(s8, 'Text 0'), ['技②　中身をきく'])
set_text(sh_of(s8, 'Text 2'), ['「確認しました」より、中身を きく。'])
set_text(sh_of(s8, 'Text 3'), ['AI「確認しました！」 ―― きみ「何を、どうやって？」'])
set_text(sh_of(s8, 'Text 5'), ['「わかっています」と 書いてあっても、そのとおりに 動くとは かぎらない。'])
set_text(sh_of(s8, 'Text 6'), [
    '跡は、おぼえている しるしには なっても、まもる しるしには ならない。',
    'それでも、きくことは 効く ――「読みました」と 言ったAIに「何行目に、何が 書いてあった？」と きいたら、読んでいなかったことが わかった。'])
# 最下部の一行（スライド7の「やってみよう」行を雛形に複製）
foot = copy.deepcopy(sh_of(S[6], 'Text 6')._element)
s8.shapes._spTree.append(foot)
f = sh_of(s8, 'Text 6')                     # 同名が二つになるため、末尾側を取り直す
cands = [sh for sh in s8.shapes if sh.name == 'Text 6']
f = cands[-1]
f._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr') \
    if False else None
nv = f._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
nv.set('name', 'Text 60'); nv.set('id', '960')
place(f, left=64, top=420, width=835, height=43)
set_text(f, ['きいた中身と、答えが 合っているか。そこまでを 一つにして、きく。'])

# ============================================================
# 3. 旧8→新9「技③ 二回きく」: 表題＋赤枠
# ============================================================
s9 = S[8]
set_text(sh_of(s9, 'Text 0'), ['技③　二回きく'])
box = copy.deepcopy(sh_of(s9, 'Shape 4')._element)
hd  = copy.deepcopy(sh_of(s9, 'Text 5')._element)
bd  = copy.deepcopy(sh_of(s9, 'Text 6')._element)
for el, nn, i in ((box, 'Shape 40', 940), (hd, 'Text 50', 950), (bd, 'Text 51', 951)):
    s9.shapes._spTree.append(el)
    nv = el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    nv.set('name', nn); nv.set('id', str(i))
b2 = [sh for sh in s9.shapes if sh.name == 'Shape 40'][0]
h2 = [sh for sh in s9.shapes if sh.name == 'Text 50'][0]
d2 = [sh for sh in s9.shapes if sh.name == 'Text 51'][0]
place(b2, top=418, height=86)
place(h2, top=428, height=30)
place(d2, top=456, height=46)
set_text(h2, ['同じ答えが 返ってきても、正しいとは かぎらない。'])
set_color(h2, RED)
set_text(d2, ['答えは 毎回、くじの引き直し。同じくじが 二回つづけて 出ることも ある。',
              'だから 二回目は、べつのくじで。同じAIに すぐ もう一回は、二回に ならない。'])

# ============================================================
# 4. 旧9→新10「技④ 逆もきく」
# ============================================================
s10 = S[9]
t0 = sh_of(s10, 'Text 0')
cur = ''.join(r.text for p in t0.text_frame.paragraphs for r in p.runs)
set_text(t0, [cur.replace('技③', '技④')])

# ============================================================
# 5. 旧16→新17「六つの やくそく」
# ============================================================
s17 = S[16]
set_text(sh_of(s17, 'Text 0'), ['六つの やくそく'])
ROWS = [('Shape 1', 'Text 2', 'Text 3'), ('Shape 4', 'Text 5', 'Text 6'),
        ('Shape 7', 'Text 8', 'Text 9'), ('Shape 10', 'Text 11', 'Text 12'),
        ('Shape 13', 'Text 14', 'Text 15')]
tree17 = s17.shapes._spTree
clone_names = ('Shape 80', 'Text 80', 'Text 81')
for nm, nn, i in zip(ROWS[0], clone_names, (980, 981, 982)):
    el = copy.deepcopy(sh_of(s17, nm)._element)
    tree17.append(el)
    nv = el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    nv.set('name', nn); nv.set('id', str(i))
ROWS.insert(1, clone_names)
TEXTS = ['すらすらを、しんじすぎない。',
         '「確認しました」より、中身をきく。',
         'だいじなことは、二回きく。',
         'よろこぶ答えほど、一回とまる。',
         'とりかえしのつかないボタンは、自分で。',
         'AIを、ただ一人の友だちに しない。']
YS = [108.0, 174.0, 240.0, 306.0, 372.0, 438.0]
for i, (grp, txt, y) in enumerate(zip(ROWS, TEXTS, YS), start=1):
    circ, num, body = [[sh for sh in s17.shapes if sh.name == n][0] for n in grp]
    place(circ, top=y + 1.5); place(num, top=y + 1.5); place(body, top=y, height=54)
    set_text(num, [str(i)]); set_text(body, [txt])

# ============================================================
# 6. ページ番号（新9以降を振り直す）
# ============================================================
for i, slide in enumerate(S, start=1):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = ''.join(r.text for p in sh.text_frame.paragraphs for r in p.runs).strip()
        if t.isdigit() and sh.left is not None and Emu(sh.left).pt > 880 and Emu(sh.top).pt > 480:
            if t != str(i): set_text(sh, [str(i)])

import os
os.makedirs(os.path.dirname(DST), exist_ok=True)
P.save(DST)
b = io.open(DST, 'rb').read()
print('保存: %s' % DST)
print('スライド枚数: %d ／ %d B ／ SHA %s' % (len(P.slides._sldIdLst), len(b),
                                       hashlib.sha256(b).hexdigest()[:16].upper()))
