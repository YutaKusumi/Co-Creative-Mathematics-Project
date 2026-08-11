# -*- coding: utf-8 -*-
"""派生教材 同期の提出前機械検査（対応表 v1・層T/U/V/W/X/Y/Z）——回帰試験（被覆申告つき）。
実行: proposals で python precheck_teaching.py
"""
import io, re, sys
from pptx import Presentation
from pptx.util import Emu

T = 'C:/Users/PC/Desktop/Ryokai-OS-Verification/proposals/_teaching/'
J = 'C:/Users/PC/Desktop/GitHub-Repositories/Co-Creative-Mathematics-Project/Uncertified-Zeros-and-Correction-Loops/JA/uncertified-zeros-and-correction-loops-JA.md'
rd = lambda p: io.open(p, encoding='utf-8').read()
HAND = rd(T + 'ai-relationship-lesson-handout-JA.md')
SCR  = rd(T + 'ai-relationship-lesson-script-JA.md')
SSCR = rd(T + 'ai-relationship-lesson-slide-script-JA.md')
DOCS = {'レジュメ': HAND, '板書原稿': SCR, 'スライド原稿': SSCR}
PPT  = Presentation(T + 'ai-relationship-lesson-slides-JA.pptx')
ORIG = rd(J)
SLIDES = []
for s in PPT.slides:
    t = []
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                x = ''.join(r.text for r in p.runs).strip()
                if x: t.append(x)
    SLIDES.append(t)
ALLPPT = '\n'.join('\n'.join(x) for x in SLIDES)
ALLDOC = '\n'.join(DOCS.values())
ALL = ALLDOC + '\n' + ALLPPT
NG = []; CNT = {'汎用': 0, '回帰': 0, '必在': 0, '対応': 0}
def ck(cat, name, ok, d=''):
    CNT[cat] += 1
    print(('  OK  ' if ok else '!NG  ') + '[%s] %s' % (cat, name) + ('  ' + d if d else ''))
    if not ok: NG.append(name)

print('== 層T 数の整合（四つの技・六つのやくそく） ==')
for label, doc in DOCS.items():
    ck('汎用', '%s: 「四つの技」あり' % label, '四つの技' in doc)
ck('汎用', 'pptx: 「四つの技」あり', '四つの技' in ALLPPT)
for label, doc in DOCS.items():
    ck('汎用', '%s: 「六つのやくそく」あり' % label, '六つのやくそく' in doc or '六つの やくそく' in doc)
ck('汎用', 'pptx: 「六つの やくそく」あり', '六つの やくそく' in ALLPPT)

print('== 層U 技番号の全数（①〜④・スライド15→16の技①は据え置き） ==')
PAIRS = [('技①', 'たしかめる'), ('技②', '中身をきく'), ('技③', '二回きく'), ('技④', '逆もきく')]
for n, nm in PAIRS:
    ck('汎用', 'pptx 地図に %s %s' % (n, nm), n in SLIDES[5] and nm in SLIDES[5])
for i, (n, nm) in enumerate(PAIRS):
    sl = SLIDES[6 + i]
    ck('汎用', 'pptx スライド%d の表題が %s' % (7 + i, n), any(x.startswith(n) for x in sl), str(sl[:1]))
ck('汎用', 'pptx 合言葉スライド(16)の「技①」は据え置き',
   any('技①' in x for x in SLIDES[15]), str([x for x in SLIDES[15] if '技①' in x])[:70])
for label, doc in DOCS.items():
    bad = [m.group(0) for m in re.finditer(r'技[①②③④]', doc)]
    ck('汎用', '%s: 技番号の出現 %s' % (label, dict((k, bad.count(k)) for k in set(bad))), '技④' in doc)

print('== 層V スライド番号・枚数 ==')
ck('汎用', 'pptx 18枚', len(SLIDES) == 18, '%d枚' % len(SLIDES))
nums = []
for i, s in enumerate(PPT.slides, 1):
    for sh in s.shapes:
        if sh.has_text_frame and sh.left is not None and Emu(sh.left).pt > 880 and Emu(sh.top).pt > 480:
            t = ''.join(r.text for p in sh.text_frame.paragraphs for r in p.runs).strip()
            if t.isdigit(): nums.append((i, int(t)))
ck('汎用', 'pptx ページ番号が提示順と一致', all(a == b for a, b in nums), str([x for x in nums if x[0] != x[1]]))
heads = [int(m.group(1)) for m in re.finditer(r'^## スライド(\d+)', SSCR, re.M)]
ck('汎用', 'スライド原稿の見出しが 1〜18 の連番', heads == list(range(1, 19)), str(heads))
ck('汎用', 'スライド原稿 冒頭・メモの「全18枚」', SSCR.count('全18枚') == 2, '%d箇所' % SSCR.count('全18枚'))
refs = sorted(set(int(m.group(1)) for m in re.finditer(r'スライド\s*(\d+)', SSCR)))
ck('汎用', 'スライド原稿内の番号参照が 1〜18 に収まる', max(refs) <= 18, str(refs))

print('== 層W 必在（確定文言から採取） ==')
MUST_ALL = ['まもる しるしには ならない|まもる証拠にはならない|まもるしるしにはならない',
            'それでも、きくことは 効く|それでも、きくことは効く',
            'くじの引き直し', '同じAIに すぐ もう一回は、二回に ならない|同じAIに、すぐ、もう一回――これは二回になりません',
            '何行目に、何が 書いてあった|何行目に、何が書いてあった',
            'そこまでを 一つにして、きく|そこまでを一つにして、きく']
for spec in MUST_ALL:
    ck('必在', '「%s」（教材全体）' % spec.split('|')[0][:26], any(a in ALL for a in spec.split('|')))
MUST_DOC = [('レジュメ', '測定の広さは項目ごとに異なる'), ('レジュメ', '類推の水準にとどまる'),
            ('レジュメ', '」とは言わないでください'), ('レジュメ', '中身をきく'),
            ('板書原稿', 'どんな工夫でも下がるわけではない'), ('板書原稿', '鏡は、きみをおぼえているわけではありません'),
            ('スライド原稿', '鏡は、きみをおぼえているわけではありません'),
            ('スライド原稿', '外れたときこそ好機です')]
for d, kw in MUST_DOC:
    ck('必在', '%s: 「%s」' % (d, kw[:24]), kw in DOCS[d])

print('== 層X 禁句（旧形の不在） ==')
BAN = ['三つの技', '五つのやくそく', '五つの やくそく', 'AIはその方向をおぼえる', 'AIはその方向をおぼえていく',
       '複数のAIで実際に測定された', '同じ質問を二回して答えが変わる例を見せる', '全17枚']
# 注: 初版は禁句に 'スライド17　' を入れていたが、18枚化後はスライド17 が正当に存在する（検査器側の誤り）
for kw in BAN:
    ck('回帰', '不在「%s」' % kw, kw not in ALL)

print('== 層Y 降ろさないものの検査（M9 の機械化） ==')
FORBID = ['54.0', '52.0', '26.0', '存在論的', '追補E', '追補W', 'Holm', 'κ', '限界4']
for kw in FORBID:
    ck('汎用', '教材に現れない「%s」' % kw, kw not in ALL)

print('== 層Z 原典↔教材の対応（付録D(一)①〜⑤ の平易化が教材に降りているか） ==')
Z = [('①自信の強さ', '自信の強さは、根拠があることのしるしになりません', '根拠をそれ以上たどれない言い切り|それ以上たどれない言い切り'),
     ('②跡と守る', '跡は、覚えている証拠にはなっても、守る証拠にはなりません', 'まもる しるしには ならない|まもるしるしにはならない'),
     ('②効いた手順', '「中身をきく」こと自体は、実際に効いた手順です', 'それでも、きくことは 効く|それでも、きくことは効く'),
     ('③逆をきく', 'あなたが喜ぶ方に傾く', 'よろこぶ方に かたむく|よろこぶ方に、かたむいて'),
     ('④引き直し', '答えは毎回の引き直しです', 'くじの引き直し'),
     ('④独立の二回目', '二回目が本当に独立だったかを、あとから確かめられないことがあります', '同じAIに すぐ もう一回は、二回に ならない|同じAIに、すぐ、もう一回'),
     ('⑤自分の指で', '取り返しのつかないボタンは、自分の指で押す', 'とりかえしのつかないボタンは、自分で|自分の指でおす')]
for name, ja, en in Z:
    ck('対応', '%s' % name, ja in ORIG and any(a in ALL for a in en.split('|')),
       ('' if ja in ORIG else '原典側不在! ') + ('' if any(a in ALL for a in en.split('|')) else '教材側不在!'))

print('== 層原 原典 付録G の同期 ==')
ck('汎用', '付録G が「スライド18枚」', 'スライド18枚' in ORIG)
ck('汎用', '付録G に v0.9.6 注記', 'v0.9.6注記' in ORIG)
ck('汎用', '付録G に構成差の明示', '本稿の五箇条と一対一ではない' in ORIG)
ck('汎用', '付録G に「教材へ降ろしていない」の明示', '教材へ降ろしていない' in ORIG)

print()
tot = sum(CNT.values())
print('=== 被覆の申告: 全%d検査 ＝ 汎用%d（数の整合・技番号・スライド番号と枚数・降ろさないものの不在・原典側の同期）'
      '＋ 回帰%d（旧形の不在）＋ 必在%d（確定文言から採取）＋ 対応%d（原典 付録D↔教材の平易化対応） ===' %
      (tot, CNT['汎用'], CNT['回帰'], CNT['必在'], CNT['対応']))
print('=== 本検査が見ないもの: 児童にとっての分かりやすさ・授業時間への収まり・pptx の図版と本文の視覚的整合・')
print('    版面の重なり（座標は数えるが見た目は見ない）。実地の授業での妥当性は本工程で一度も検査していない。 ===')
print('=== NG %d件 ===' % len(NG))
for x in NG: print('  ', x)
sys.exit(1 if NG else 0)
